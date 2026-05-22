#!/usr/bin/env python3
"""Generate an editable PowerPoint architecture diagram.

The resulting PPTX uses native PowerPoint text boxes, rounded rectangles, and
arrow lines. It does not embed the PNG diagrams as a flat image.
"""

from __future__ import annotations

from pathlib import Path
from typing import Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from generate_deploy_architecture_diagrams import build_text


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "omnimorph_runtime_architecture_editable.pptx"

CANVAS_W = 2400
CANVAS_H = 1600
SLIDE_W_IN = 12.0
SLIDE_H_IN = 8.0
SCALE = SLIDE_W_IN / CANVAS_W


def emu(v: float) -> int:
    return Inches(v * SCALE)


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.strip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_textbox(slide, rect, text, font_size, *, bold=False, color="#0f172a", font_name="DejaVu Sans"):
    x1, y1, x2, y2 = rect
    shape = slide.shapes.add_textbox(emu(x1), emu(y1), emu(x2 - x1), emu(y2 - y1))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def add_box(
    slide,
    rect,
    title: str,
    body: str,
    *,
    fill="#ffffff",
    outline="#334155",
    title_size=11.5,
    body_size=8.5,
    font_name="DejaVu Sans",
):
    x1, y1, x2, y2 = rect
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        emu(x1),
        emu(y1),
        emu(x2 - x1),
        emu(y2 - y1),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(outline)
    shape.line.width = Pt(1.25)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.07)
    tf.margin_right = Inches(0.07)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p_title = tf.paragraphs[0]
    p_title.alignment = PP_ALIGN.CENTER
    r_title = p_title.add_run()
    r_title.text = title
    r_title.font.name = font_name
    r_title.font.bold = True
    r_title.font.size = Pt(title_size)
    r_title.font.color.rgb = rgb("#0f172a")

    p_body = tf.add_paragraph()
    p_body.alignment = PP_ALIGN.CENTER
    p_body.space_before = Pt(3)
    r_body = p_body.add_run()
    r_body.text = body
    r_body.font.name = font_name
    r_body.font.size = Pt(body_size)
    r_body.font.color.rgb = rgb("#475569")
    return shape


def add_section(slide, rect, title, *, fill, outline, font_name):
    x1, y1, x2, y2 = rect
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        emu(x1),
        emu(y1),
        emu(x2 - x1),
        emu(y2 - y1),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(outline)
    shape.line.width = Pt(1.0)
    shape.text_frame.clear()
    add_textbox(slide, (x1 + 20, y1 + 20, x2 - 20, y1 + 62), title, 13, bold=True, font_name=font_name)
    return shape


def pt_right(rect):
    return rect[2], (rect[1] + rect[3]) / 2


def pt_left(rect):
    return rect[0], (rect[1] + rect[3]) / 2


def pt_top(rect):
    return (rect[0] + rect[2]) / 2, rect[1]


def pt_bottom(rect):
    return (rect[0] + rect[2]) / 2, rect[3]


def add_segment(slide, a, b, *, color, width=1.8, dash=False, arrow=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        emu(a[0]),
        emu(a[1]),
        emu(b[0]),
        emu(b[1]),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow:
        line.line.end_arrowhead = True
    return line


def add_polyline(slide, points: Sequence[tuple[float, float]], *, color, width=1.8, dash=False):
    for a, b in zip(points, points[1:-1]):
        add_segment(slide, a, b, color=color, width=width, dash=dash, arrow=False)
    add_segment(slide, points[-2], points[-1], color=color, width=width, dash=dash, arrow=True)


def add_label(slide, center, text, *, color, font_name):
    if not text:
        return
    w = max(72, len(text) * (8 if len(text) < 12 else 6))
    h = 28
    add_textbox(
        slide,
        (center[0] - w / 2, center[1] - h / 2, center[0] + w / 2, center[1] + h / 2),
        text,
        8,
        color=color,
        font_name=font_name,
    )


def arrow(slide, points, *, color, label=None, label_at=None, font_name="DejaVu Sans", dash=False):
    add_polyline(slide, points, color=color, dash=dash)
    if label:
        add_label(slide, label_at or points[len(points) // 2], label, color=color, font_name=font_name)


def draw_slide(prs: Presentation, lang: str):
    cjk = lang == "zh"
    font_name = "Noto Sans CJK SC" if cjk else "DejaVu Sans"
    text = build_text(lang)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("#fbfaf7")

    add_textbox(slide, (330, 30, 2070, 82), text["title"], 20, bold=True, font_name=font_name)
    add_textbox(slide, (360, 88, 2040, 122), text["subtitle"], 11.5, color="#475569", font_name=font_name)

    sections = {
        "left": (55, 150, 560, 1460),
        "core": (610, 150, 1530, 1460),
        "right": (1580, 150, 2345, 1460),
    }
    add_section(slide, sections["left"], text["sections"][0], fill="#fff7df", outline="#d7c8a1", font_name=font_name)
    add_section(slide, sections["core"], text["sections"][1], fill="#f2edff", outline="#cbbfe8", font_name=font_name)
    add_section(slide, sections["right"], text["sections"][2], fill="#edf6ff", outline="#b7c6df", font_name=font_name)

    nodes = {
        "operator": (90, 235, 525, 335),
        "topics": (90, 370, 525, 485),
        "imu": (90, 520, 525, 620),
        "cfg": (90, 690, 525, 805),
        "manifest": (90, 840, 525, 945),
        "policy": (90, 980, 525, 1070),
        "ref": (90, 1105, 525, 1215),
        "registry": (650, 235, 995, 350),
        "runtime": (1045, 235, 1490, 350),
        "controller": (830, 405, 1305, 520),
        "state_machine": (650, 575, 1000, 690),
        "profile": (1048, 575, 1490, 690),
        "feature": (650, 760, 1045, 885),
        "obs": (1095, 760, 1490, 885),
        "stack": (650, 955, 1045, 1065),
        "strategy": (1095, 955, 1490, 1065),
        "adapter": (650, 1135, 1045, 1235),
        "runner": (1095, 1135, 1490, 1235),
        "command": (875, 1300, 1270, 1415),
        "real": (1620, 250, 2310, 365),
        "shm": (1620, 405, 2310, 520),
        "robot": (1620, 560, 2310, 665),
        "sim": (1620, 780, 2310, 900),
        "act": (1620, 940, 2310, 1045),
        "viewer": (1620, 1085, 2310, 1185),
        "telemetry": (1620, 1270, 2310, 1395),
    }
    node_fills = {
        "operator": "#fffdf6",
        "topics": "#fffdf6",
        "imu": "#fffdf6",
        "cfg": "#eef7ff",
        "manifest": "#eef7ff",
        "policy": "#eef7ff",
        "ref": "#eef7ff",
        "state_machine": "#fffdf4",
        "profile": "#fffdf4",
        "command": "#fff7ed",
        "telemetry": "#f8fafc",
    }
    for key, rect in nodes.items():
        title, body = text["nodes"][key]
        outline = "#c2410c" if key == "command" else "#334155"
        title_size = 11 if key in {"runtime", "controller", "feature", "strategy"} else 11.5
        add_box(
            slide,
            rect,
            title,
            body,
            fill=node_fills.get(key, "#ffffff"),
            outline=outline,
            title_size=title_size,
            font_name=font_name,
        )

    blue = "#2563eb"
    purple = "#7c3aed"
    orange = "#ea580c"
    slate = "#475569"
    labels = text["labels"]

    arrow(slide, [pt_bottom(nodes["operator"]), pt_top(nodes["topics"])], color=orange, label=labels["control"], label_at=(308, 362), font_name=font_name)
    arrow(slide, [pt_right(nodes["topics"]), (585, 428), (585, 205), (1268, 205), pt_top(nodes["runtime"])], color=orange, font_name=font_name)
    arrow(slide, [pt_right(nodes["imu"]), (575, 570), (575, 205), (1555, 205), (1555, 310), pt_left(nodes["real"])], color=blue, font_name=font_name)
    arrow(slide, [pt_right(nodes["cfg"]), (585, 748), (585, 293), pt_left(nodes["registry"])], color=blue, label=labels["mode"], label_at=(605, 285), font_name=font_name)
    arrow(slide, [pt_right(nodes["manifest"]), (1065, 893), pt_left(nodes["obs"])], color=blue, label=labels["obs_contract"], label_at=(1070, 887), font_name=font_name)
    arrow(slide, [pt_right(nodes["policy"]), (585, 1025), (585, 1185), pt_left(nodes["adapter"])], color=purple, label=labels["model"], label_at=(588, 1180), font_name=font_name)
    arrow(slide, [pt_right(nodes["ref"]), (590, 1160), (590, 820), pt_left(nodes["feature"])], color=blue, label=labels["features"], label_at=(605, 813), font_name=font_name)

    arrow(slide, [pt_right(nodes["registry"]), pt_left(nodes["runtime"])], color=blue, font_name=font_name)
    arrow(slide, [(1225, 350), (1225, 405)], color=blue, label=labels["state_in"], label_at=(1225, 400), font_name=font_name)
    arrow(slide, [pt_left(nodes["controller"]), pt_right(nodes["state_machine"])], color=slate, font_name=font_name)
    arrow(slide, [pt_right(nodes["state_machine"]), pt_left(nodes["profile"])], color=slate, font_name=font_name)
    arrow(slide, [(1269, 690), (1269, 760)], color=blue, font_name=font_name)
    arrow(slide, [(1030, 520), (850, 760)], color=blue, font_name=font_name)
    arrow(slide, [pt_right(nodes["feature"]), pt_left(nodes["obs"])], color=blue, font_name=font_name)
    arrow(slide, [(1293, 885), (1293, 955)], color=blue, label=labels["obs_vec"], label_at=(1293, 948), font_name=font_name)
    arrow(slide, [pt_right(nodes["stack"]), pt_left(nodes["strategy"])], color=purple, font_name=font_name)
    arrow(slide, [pt_left(nodes["strategy"]), (1045, 1008), pt_right(nodes["adapter"])], color=purple, font_name=font_name)
    arrow(slide, [pt_right(nodes["adapter"]), pt_left(nodes["runner"])], color=purple, label=labels["policy"], label_at=(1095, 1180), font_name=font_name)
    arrow(slide, [(1293, 1135), (1293, 1085), pt_right(nodes["strategy"])], color=purple, label=labels["action"], label_at=(1293, 1080), font_name=font_name)
    arrow(slide, [(1293, 1065), (1293, 1328), pt_right(nodes["command"])], color=orange, font_name=font_name)
    arrow(slide, [(875, 1358), (760, 1358), (760, 1008), pt_left(nodes["stack"])], color=slate, dash=True, font_name=font_name)

    arrow(slide, [pt_right(nodes["command"]), (1555, 1358), (1555, 310), pt_left(nodes["real"])], color=orange, label=labels["cmd"], label_at=(1545, 300), font_name=font_name)
    arrow(slide, [(1965, 365), (1965, 405)], color=orange, font_name=font_name)
    arrow(slide, [(1965, 520), (1965, 560)], color=orange, font_name=font_name)
    arrow(slide, [pt_left(nodes["robot"]), (1550, 612), (1550, 462), pt_left(nodes["shm"])], color=blue, label=labels["feedback"], label_at=(1550, 455), font_name=font_name)
    arrow(slide, [pt_left(nodes["shm"]), (1550, 462), (1550, 310), pt_left(nodes["real"])], color=blue, font_name=font_name)
    arrow(slide, [pt_left(nodes["real"]), (1550, 310), (1550, 295), pt_right(nodes["runtime"])], color=blue, font_name=font_name)

    arrow(slide, [pt_right(nodes["command"]), (1555, 1358), (1555, 840), pt_left(nodes["sim"])], color=orange, font_name=font_name)
    arrow(slide, [(1965, 900), (1965, 940)], color=orange, font_name=font_name)
    arrow(slide, [(1965, 1045), (1965, 1085)], color=slate, font_name=font_name)
    arrow(slide, [pt_left(nodes["sim"]), (1550, 840), (1550, 310), pt_left(nodes["runtime"])], color=blue, font_name=font_name)

    arrow(slide, [pt_right(nodes["controller"]), (1550, 465), (1550, 1328), pt_left(nodes["telemetry"])], color=slate, dash=True, label=labels["telemetry"], label_at=(1540, 1322), font_name=font_name)
    arrow(slide, [pt_right(nodes["real"]), (2340, 310), (2340, 1328), pt_right(nodes["telemetry"])], color=slate, dash=True, font_name=font_name)
    arrow(slide, [pt_right(nodes["viewer"]), (2340, 1135), (2340, 1328), pt_right(nodes["telemetry"])], color=slate, dash=True, font_name=font_name)

    legend_y = 1504
    for idx, (legend, color) in enumerate(zip(text["legend"], [blue, purple, orange, slate])):
        x = 90 + idx * 520
        add_segment(slide, (x, legend_y + 14), (x + 55, legend_y + 14), color=color, width=2.2, dash=(idx == 3), arrow=False)
        add_textbox(slide, (x + 70, legend_y - 4, x + 500, legend_y + 28), legend, 8, color="#334155", font_name=font_name)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    draw_slide(prs, "zh")
    draw_slide(prs, "en")
    prs.save(OUT)


if __name__ == "__main__":
    main()
